"""
Generate IpsosKE QC Dashboard User Guide in DOCX, PDF, and PPTX formats.
Run from the project root:  python generate_guide.py
"""

import os

NAVY  = (31, 43, 108)
TEAL  = (0, 181, 173)
ORANGE = (255, 112, 67)
WHITE = (255, 255, 255)
GRAY  = (245, 245, 245)
DARK  = (33, 33, 33)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ─────────────────────────────────────────────────────────────────────────────
# Content definition
# ─────────────────────────────────────────────────────────────────────────────

SECTIONS = [
    {
        "title": "1. Getting Started",
        "body": (
            "Open a terminal, navigate to the qc_app folder inside the project directory, "
            "then run the following command:\n\n"
            "    python -m streamlit run app.py\n\n"
            "The dashboard will open automatically in your default web browser at "
            "http://localhost:8501."
        ),
    },
    {
        "title": "2. Logging In",
        "body": (
            "Option A — Guest (Prototype mode)\n"
            "Click 'Continue as Guest' on the login page. You receive full QC Executive "
            "access with no account required. A yellow GUEST badge appears in the sidebar.\n\n"
            "Option B — Account Login\n"
            "Enter your email and password, then click Login.\n\n"
            "Default admin account (seeded on first run):\n"
            "  • Email:    admin@ipsos.com\n"
            "  • Password: admin1234\n\n"
            "Self-registration creates an 'Other' account (dashboard view only). "
            "An admin must raise the role before the user can upload data or drill into projects."
        ),
    },
    {
        "title": "3. User Roles & Permissions",
        "body": (
            "QC Executive — Full access: create users, projects, upload data, drill into all projects.\n"
            "Operations Manager — Same as QC Executive.\n"
            "QC Officer — Upload data and drill into assigned projects only.\n"
            "Project Manager — View and drill into assigned projects. No upload rights.\n"
            "Researcher — View and drill into assigned projects. No upload rights.\n"
            "Management — View all projects summary. No drill-down.\n"
            "Other — Summary dashboard view only. No drill-down."
        ),
    },
    {
        "title": "4. Dashboard (Home Page)",
        "body": (
            "After login you land on the Dashboard showing:\n\n"
            "• KPI Strip — Active projects, Total Sample Target, Total Approved interviews, "
            "Average Back-check Rate, Total Flagged records.\n"
            "• Completion Chart — Grouped bar chart: Target vs Approved per project.\n"
            "• Project Table — Every project with progress bar, BC rate, Listen-in rate, "
            "Flagged count, and a View Details button (for roles with drill-down rights).\n"
            "• Risk & Alert Status — Auto-flags projects nearing their end date, below "
            "back-check target, or with a high flagged rate.\n"
            "• Recent Activity — Last 10 uploads across all projects.\n\n"
            "Filters at the top let you narrow by Status (Active / Completed / Paused) "
            "and by project End Date range."
        ),
    },
    {
        "title": "5. Admin Panel",
        "body": (
            "Accessible from the sidebar to QC Executive and Operations Manager only.\n\n"
            "5.1 User Management\n"
            "• View all users in the table.\n"
            "• Change a user's role: select the user, pick a new role, click Update Role.\n"
            "• Activate / Deactivate: select user, choose action, click Apply.\n"
            "• Add a new user: fill in Name, Email, Role, Temporary Password, click Create User.\n\n"
            "5.2 Project Management\n"
            "Create a project:\n"
            "  1. Go to Admin -> Project Management.\n"
            "  2. Fill in: Project Name, Client, Sample Target, Start Date, End Date.\n"
            "  3. Set targets with sliders: Back-check % (default 20%), "
            "Listen-in % (default 10%), Accompaniment % (default 20%).\n"
            "  4. Click Create Project.\n\n"
            "Edit a project:\n"
            "  1. Select the project in 'Edit Project'.\n"
            "  2. All fields are pre-filled — change what you need.\n"
            "  3. Update the Status (Active / Paused / Completed).\n"
            "  4. Click Save Changes.\n\n"
            "5.3 Project Assignments\n"
            "  1. Select a project.\n"
            "  2. Current assignees are listed — click Remove to unassign.\n"
            "  3. Multiselect users to add, then click Assign Selected.\n\n"
            "5.4 Upload History\n"
            "Shows every file uploaded across all projects. Admins can delete any upload — "
            "this removes both the log entry and all associated data rows."
        ),
    },
    {
        "title": "6. Project Detail",
        "body": (
            "Click View Details on any project row in the Dashboard. "
            "You must be assigned to the project (or be QC Executive / Operations Manager).\n\n"
            "Header Banner shows at a glance:\n"
            "  • Project name and client\n"
            "  • Sample Target vs Achieved (±overflow in colour)\n"
            "  • Back-check rate vs target ([OK] or [!])\n"
            "  • Listen-in rate vs target ([OK] or [!])\n"
            "  • Project status badge\n\n"
            "Seven sub-tabs are available:\n"
            "  Quality Report | Back-check Report | Cancelled Interviews | "
            "Performance Report | Timing Report | Listen-in | Wave Comparison"
        ),
    },
    {
        "title": "7. Quality Report Tab",
        "body": (
            "Uploading Data:\n"
            "  1. Expand 'Upload New Data' (opens automatically when no data exists).\n"
            "  2. Optionally enter a Wave / Period label (e.g. 'Wave 1', 'April 2025').\n"
            "  3. Click 'Download Template (.xlsx)' for the correct column format.\n"
            "  4. Upload your iField Excel export.\n"
            "  5. Review the Column Mapping UI:\n"
            "     [Mapped] = required column mapped\n"
            "     [Missing] = required column missing — assign manually\n"
            "     [Optional] = optional column\n"
            "  6. Select any extra project-specific columns to retain.\n"
            "  7. Click Confirm & Save Data.\n\n"
            "Required columns: INSTANCE_ID, INTERVIEWER_ID, INTERVIEW_DATE, "
            "INTERVIEW_START_TIME, INTERVIEW_END_TIME, DURATION_MINUTES, DURATION_FLAG, "
            "STRAIGHT_LINING, LONG_PAUSE, REGION, LOCATION, SAMPLE_POINT_ID, APPROVAL_STATUS.\n\n"
            "Optional columns (auto-detected): GPS_STATUS, PHONE_PRESENT, AUDIO_PRESENT, "
            "DURATION_VALIDATION, QC_COMMENTS.\n\n"
            "What the tab shows after upload:\n"
            "  • KPIs: Total Submitted, Approved, Pending, Cancelled, Flagged, "
            "Completion %, Avg/Max/Min Duration, Error Rate\n"
            "  • Back-checks & Accompaniments block (live count vs target)\n"
            "  • Charts: Approval donut, Error type donut, Errors by interviewer, "
            "Productivity trend, Duration by interviewer, Duration by region, "
            "Status by region\n"
            "  • Filters above charts: Interviewer, Region, Status\n"
            "  • Quality Queries: auto-detected issues list\n"
            "  • Mitigation Measures: standard bullets always shown\n"
            "  • Export: Excel (.xlsx) and SAV/SPSS (.sav)"
        ),
    },
    {
        "title": "8. Back-check Report Tab",
        "body": (
            "Uploading: Use the template with columns BC_INSTANCE_ID, ORIGINAL_INSTANCE_ID, "
            "INTERVIEW_STATUS, REGION, LOCATION, SAMPLE_POINT_ID, BACKCHECKER_ID, "
            "INTERVIEWER_ID, SCRIPT_NAME, INTERVIEW_DATE, BACKCHECK_DATE, "
            "ERROR_01 through ERROR_13. Error columns are binary (1 = error, 0 = none). "
            "iField column names auto-map.\n\n"
            "What it shows:\n"
            "  • KPIs: Approved interviews, Back-checks Done, Effective Back-checks, "
            "BC Rate %, Target %\n"
            "  • Gauge chart: Back-check rate vs target\n"
            "  • Error type donut: which of the 13 error codes were triggered\n"
            "  • Errors by interviewer (stacked bar)\n"
            "  • Errors by region (stacked bar)\n"
            "  • Export: Excel + SAV\n\n"
            "Error codes:\n"
            "  01 Totally fraudulent interview\n"
            "  02 Different respondent name/address\n"
            "  03 Mismatched quotas (age/gender)\n"
            "  04 Wrong usership/recruitment criteria\n"
            "  05 Wrong telephone number\n"
            "  06 Unattainable/invalid telephone\n"
            "  07 Engaged/Answer machine\n"
            "  08 Refused/unable to reach respondent\n"
            "  09 Respondent doesn't remember interview\n"
            "  10 Back-check abandoned/incomplete\n"
            "  11 Interview done on wrong mode\n"
            "  12 Voice recording permission issue\n"
            "  13 Respondent already participated recently"
        ),
    },
    {
        "title": "9. Cancelled Interviews Tab",
        "body": (
            "Uploading: Use the template with columns INSTANCE_ID, REGION, INTERVIEWER_ID, "
            "INTERVIEW_DATE, INTERVIEW_LENGTH, ACTIVE_LENGTH, GAP_TO_LAST, QF_A through QF_F, "
            "INTERVIEWER_PERFORMANCE, BACKCHECK_RESULT_TELEPHONE/F2F/INDEPENDENT.\n\n"
            "What it shows:\n"
            "  • KPIs: Total Cancellations, Cancellation Rate %, Avg Duration, Avg Gap to Last\n"
            "  • Cancellations by Interviewer (horizontal bar)\n"
            "  • Cancellations by Region (bar)\n"
            "  • Quality Flags Distribution (QF A–F) — donut chart\n"
            "  • Gap to Last Interview — histogram of time gaps between consecutive interviews\n"
            "  • Export: Excel + SAV"
        ),
    },
    {
        "title": "10. Performance Report Tab",
        "body": (
            "Uploading: iField Performance Report has two header rows — the system "
            "skips them automatically. Template columns: INTERVIEWER_ID, REGION, "
            "FIRST_INTERVIEW, LAST_INTERVIEW, INTERVIEW_COMPLETES, FOLLOWUP_COMPLETES, "
            "ECS_COMPLETES, WORK_SUMMARY, ACCOMPANIMENTS, CANCELLED_INTERVIEWS, "
            "BACKCHECK_TELEPHONE_CREATED, BACKCHECK_F2F_CREATED, BACKCHECK_F2F_INFIELD, "
            "BACKCHECK_TOTAL, BACKCHECK_COMPLETED.\n\n"
            "What it shows:\n"
            "  • KPIs: Total Interviewers, Interview Completes, Accompaniments, "
            "Accompaniment Rate %, BC Completed\n"
            "  • Gauge charts: Accompaniment rate vs target, BC completion rate vs target\n"
            "  • Completes by Interviewer (sorted bar)\n"
            "  • Completes by Region (bar) + Accompaniment donut\n"
            "  • Back-check Activity Summary (bar): telephone, F2F created, F2F in-field, completed\n"
            "  • Team Retention: interviewers active across multiple projects\n"
            "  • Export: Excel + SAV"
        ),
    },
    {
        "title": "11. Timing Report Tab",
        "body": (
            "Uploading: Template columns: INSTANCE_ID, INTERVIEWER_ID, REGION, "
            "INTERVIEW_DATE, DURATION_MINUTES. If no Timing Report is uploaded, "
            "the tab automatically falls back to duration data from the Quality Report.\n\n"
            "What it shows:\n"
            "  • KPIs: Average, Median, Maximum, Minimum duration (minutes), "
            "Interviews Below 50% of Average (flagged)\n"
            "  • Duration Distribution histogram with Average line and 50%-avg threshold\n"
            "  • Avg Duration by Interviewer (bar)\n"
            "  • Avg Duration by Region (bar)\n"
            "  • Duration Trend Over Time per interviewer (line) with threshold line\n"
            "  • Duration vs Productivity scatter (if Performance Report data exists)\n"
            "  • Export: Excel + SAV"
        ),
    },
    {
        "title": "12. Listen-in Tab",
        "body": (
            "Logging a single session:\n"
            "  1. Click 'Log Single Session'.\n"
            "  2. Select Interviewer ID (pre-populated from Quality Report) or type manually.\n"
            "  3. Fill in: Listen Date, Listen Type (Telephone / F2F Accompaniment / "
            "Audio Playback), Result (Pass / Fail / Partial).\n"
            "  4. Note any Issues and Actions Taken.\n"
            "  5. Click Save Session.\n\n"
            "Batch upload:\n"
            "  1. Click 'Batch Upload'.\n"
            "  2. Enter a Wave / Period label.\n"
            "  3. Download the template (columns: INSTANCE_ID, INTERVIEWER_ID, REGION, "
            "LISTEN_DATE, LISTEN_TYPE, RESULT, ISSUES_NOTED, ACTION_TAKEN).\n"
            "  4. Upload the filled file and click Save Batch.\n\n"
            "What it shows:\n"
            "  • KPIs: Approved Interviews, Listen-in Sessions, Listen-in Rate %, Target %, Pass Rate %\n"
            "  • Gauge chart: rate vs target\n"
            "  • Result Breakdown donut (Pass / Fail / Partial)\n"
            "  • Listen Type Breakdown donut\n"
            "  • Sessions by Interviewer (bar)\n"
            "  • Daily Trend line (split by type)\n"
            "  • Sessions with Issues table\n"
            "  • Full session log with delete option\n"
            "  • Export: Excel + SAV"
        ),
    },
    {
        "title": "13. Wave Comparison Tab",
        "body": (
            "This tab activates once at least one upload on any tab has a Wave / Period label set.\n\n"
            "To populate it: when uploading any report, fill in the 'Wave / Period label' field "
            "before saving. Use consistent names across all uploads in the same wave "
            "(e.g. always 'Wave 1').\n\n"
            "What it shows:\n"
            "  • KPI row: Waves tracked, Total Approved, Avg Error Rate, Avg BC Rate, Avg Duration\n"
            "  • Error Rate Trend: Estimated (rolling avg of prior waves) vs Actual + "
            "Variance bars (orange = worse, teal = better)\n"
            "  • Back-check Rate by Wave (bar, colour-coded vs target)\n"
            "  • Listen-in Rate by Wave (bar vs target)\n"
            "  • Approved Interviews + Productivity (dual-axis: bar + line)\n"
            "  • Avg Duration by Wave (area line)\n"
            "  • Wave Summary Table: all metrics in one exportable view"
        ),
    },
    {
        "title": "14. Exporting Data",
        "body": (
            "Every tab has Export to Excel (.xlsx) and Export to SAV/SPSS (.sav) buttons "
            "at the bottom of the data table.\n\n"
            "• Excel files use Ipsos Navy header formatting and are suitable for most stakeholders.\n"
            "• SAV files are compatible with SPSS and IBM Statistics for researchers and DP teams "
            "working offline.\n\n"
            "The Admin -> Upload History tab allows downloading the full upload log as Excel."
        ),
    },
    {
        "title": "15. Typical End-to-End Workflow",
        "body": (
            "Step 1:  Admin creates the project (Admin -> Project Management -> Add New Project)\n"
            "Step 2:  Admin assigns QC Officer and Project Manager to the project\n"
            "Step 3:  QC Officer logs in -> Dashboard -> View Details on their project\n"
            "Step 4:  QC Officer uploads Quality Report data (set wave label if applicable)\n"
            "Step 5:  QC Officer uploads Back-check Report\n"
            "Step 6:  QC Officer uploads Performance Report\n"
            "Step 7:  QC Officer logs Listen-in sessions (manual or batch)\n"
            "Step 8:  Project Manager / Management log in to view progress on Dashboard\n"
            "Step 9:  Repeat Steps 4–7 for each subsequent wave using the same wave label\n"
            "Step 10: View Wave Comparison tab to track trends and error rate changes across waves"
        ),
    },
]


# ─────────────────────────────────────────────────────────────────────────────
# DOCX
# ─────────────────────────────────────────────────────────────────────────────

def make_docx():
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    def set_cell_bg(cell, hex_color):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_color)
        tcPr.append(shd)

    # Cover / title block
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("IpsosKE QC Dashboard")
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*NAVY)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = p2.add_run("User Guide")
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(*TEAL)

    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run3 = p3.add_run("Ipsos Kenya — Internal Use Only")
    run3.font.size = Pt(11)
    run3.font.color.rgb = RGBColor(100, 100, 100)

    doc.add_paragraph()

    # Sections
    for sec in SECTIONS:
        # Section heading
        h = doc.add_heading(level=1)
        h.clear()
        run = h.add_run(sec["title"])
        run.font.color.rgb = RGBColor(*NAVY)
        run.font.size = Pt(14)
        run.font.bold = True

        # Body text — handle multi-paragraph blocks
        for para_text in sec["body"].split("\n\n"):
            lines = para_text.strip().split("\n")
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(2)
                if line.startswith("  ") or line.startswith("    "):
                    p.paragraph_format.left_indent = Inches(0.3)
                run = p.add_run(line.lstrip())
                run.font.size = Pt(10.5)
                run.font.color.rgb = RGBColor(*DARK)

        doc.add_paragraph()

    path = os.path.join(OUT_DIR, "QC_Dashboard_User_Guide.docx")
    doc.save(path)
    print(f"DOCX saved -> {path}")


# ─────────────────────────────────────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────────────────────────────────────

def make_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, HRFlowable, PageBreak
    )
    from reportlab.lib.enums import TA_CENTER, TA_LEFT

    path = os.path.join(OUT_DIR, "QC_Dashboard_User_Guide.pdf")

    navy_rl  = colors.Color(NAVY[0]/255,  NAVY[1]/255,  NAVY[2]/255)
    teal_rl  = colors.Color(TEAL[0]/255,  TEAL[1]/255,  TEAL[2]/255)
    dark_rl  = colors.Color(DARK[0]/255,  DARK[1]/255,  DARK[2]/255)
    gray_rl  = colors.Color(0.96, 0.96, 0.96)

    doc = SimpleDocTemplate(
        path, pagesize=A4,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        topMargin=2*cm, bottomMargin=2*cm,
    )

    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "Title", fontSize=26, textColor=navy_rl,
        fontName="Helvetica-Bold", alignment=TA_CENTER, spaceAfter=6,
    )
    subtitle_style = ParagraphStyle(
        "Subtitle", fontSize=16, textColor=teal_rl,
        fontName="Helvetica", alignment=TA_CENTER, spaceAfter=4,
    )
    meta_style = ParagraphStyle(
        "Meta", fontSize=10, textColor=colors.grey,
        fontName="Helvetica", alignment=TA_CENTER, spaceAfter=20,
    )
    h1_style = ParagraphStyle(
        "H1", fontSize=13, textColor=navy_rl,
        fontName="Helvetica-Bold", spaceBefore=14, spaceAfter=6,
        borderPad=4,
    )
    body_style = ParagraphStyle(
        "Body", fontSize=10, textColor=dark_rl,
        fontName="Helvetica", spaceAfter=3, leading=15,
    )
    indent_style = ParagraphStyle(
        "Indent", fontSize=10, textColor=dark_rl,
        fontName="Helvetica", spaceAfter=2, leading=14,
        leftIndent=20,
    )

    story = []

    # Cover
    story.append(Spacer(1, 1.5*cm))
    story.append(Paragraph("IpsosKE QC Dashboard", title_style))
    story.append(Paragraph("User Guide", subtitle_style))
    story.append(Paragraph("Ipsos Kenya — Internal Use Only", meta_style))
    story.append(HRFlowable(width="100%", thickness=2, color=teal_rl, spaceAfter=20))

    for sec in SECTIONS:
        story.append(Paragraph(sec["title"], h1_style))
        story.append(HRFlowable(width="100%", thickness=0.5, color=colors.Color(0.85,0.85,0.85), spaceAfter=6))

        for para_text in sec["body"].split("\n\n"):
            lines = para_text.strip().split("\n")
            for line in lines:
                raw = line.rstrip()
                if not raw:
                    continue
                is_indented = raw.startswith("  ") or raw.startswith("    ")
                clean = raw.lstrip()
                # Escape ampersands for ReportLab
                clean = clean.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                if is_indented:
                    story.append(Paragraph(clean, indent_style))
                else:
                    story.append(Paragraph(clean, body_style))
            story.append(Spacer(1, 4))

        story.append(Spacer(1, 8))

    doc.build(story)
    print(f"PDF saved  -> {path}")


# ─────────────────────────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────────────────────────

def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as PptxRGB
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def rgb(t):
        return PptxRGB(*t)

    def add_rect(slide, left, top, width, height, fill_color):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
        shape.line.fill.background()
        return shape

    def add_text_box(slide, text, left, top, width, height,
                     font_size=18, bold=False, color=DARK,
                     align=PP_ALIGN.LEFT, word_wrap=True):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        txBox.word_wrap = word_wrap
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return txBox

    # ── Title slide ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 5.8, 13.33, 0.25, TEAL)

    add_text_box(slide, "IpsosKE QC Dashboard",
                 1, 1.5, 11, 1.5, font_size=40, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "User Guide",
                 1, 3.2, 11, 1, font_size=26, bold=False, color=TEAL,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "Ipsos Kenya  ·  Internal Use Only",
                 1, 4.4, 11, 0.6, font_size=14, color=(180, 190, 220),
                 align=PP_ALIGN.CENTER)

    # ── Content slides ───────────────────────────────────────────────────────
    for sec in SECTIONS:
        slide = prs.slides.add_slide(blank_layout)

        # Header bar
        add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
        add_rect(slide, 0, 1.1, 13.33, 0.05, TEAL)

        # Title
        add_text_box(slide, sec["title"],
                     0.3, 0.15, 12.7, 0.85,
                     font_size=20, bold=True, color=WHITE)

        # Body — truncate to fit slide
        body_lines = []
        for para in sec["body"].split("\n\n"):
            for line in para.strip().split("\n"):
                body_lines.append(line.rstrip())
            body_lines.append("")

        # Use a text box for body
        txBox = slide.shapes.add_textbox(
            Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.9)
        )
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True

        first = True
        for line in body_lines[:55]:  # cap lines per slide
            raw = line.rstrip()
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            is_indent = raw.startswith("  ") or raw.startswith("    ")
            p.level = 1 if is_indent else 0
            run = p.add_run()
            run.text = raw.lstrip()
            run.font.size = Pt(10.5)
            run.font.color.rgb = rgb(DARK)

        # Footer
        add_rect(slide, 0, 7.2, 13.33, 0.3, (240, 240, 245))
        add_text_box(slide, "IpsosKE QC Dashboard  ·  Internal Use Only",
                     0.3, 7.2, 12.7, 0.3,
                     font_size=8, color=(120, 120, 130), align=PP_ALIGN.CENTER)

    path = os.path.join(OUT_DIR, "QC_Dashboard_User_Guide.pptx")
    prs.save(path)
    print(f"PPTX saved -> {path}")


# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    make_docx()
    make_pdf()
    make_pptx()
    print("\nAll three files saved to:", OUT_DIR)
